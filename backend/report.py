"""
문서 내보내기(리포트 생성) 파이프라인.

분석 결과를 docx / pptx / xlsx / txt 파일로 출력한다. NotebookLM 식 "기능 = 프롬프트
템플릿" 발상을 따르되, 로컬 소형 모델(Gemma) 한계를 감안해 **출력 형식을 JSON 스키마로
강제**하고(format=json) 실제 파일 조립은 코드가 결정론적으로 수행한다.

흐름:
  build_report(질문, 문서들) →  RAG 검색(rag.search 재사용)
                            →  Gemma 에 "리포트 JSON 으로만 답하라" 강제
                            →  검증된 content 모델(dict) 반환
  to_docx / to_pptx / to_xlsx / to_txt(content, path)  →  포맷별 렌더러

content 모델 (포맷 무관 공통):
  {
    "title": str,
    "summary": str,
    "sections": [{"heading": str, "paragraphs": [str, ...]}, ...],
    "table": {"columns": [str, ...], "rows": [[셀, ...], ...]} | None,
    "chart": {"type": "bar|line|pie", "title": str,
              "categories": [str, ...],
              "series": [{"name": str, "values": [num, ...]}, ...]} | None,
    "sources": [{"name", "page", "doc_id", ...}, ...],
    "meta": {"question", "model", "generated"},
  }

차트는 pptx/xlsx 에서 **네이티브 차트**(편집 가능)로 들어간다. matplotlib 등 외부
렌더러·폰트 의존이 없어 한글 깨짐 문제가 없고 번들이 가볍다. docx/txt 는 데이터 표로 싣는다.
"""

import json
import os
import re
import time

import requests

from backend import rag
from backend.ollama import OLLAMA_HOST

# 내보내기 지원 포맷 → 확장자
FORMATS = {
    "docx": ".docx",
    "pptx": ".pptx",
    "xlsx": ".xlsx",
    "txt": ".txt",
}

# 리포트는 표/추세를 다루므로 일반 Q&A(top_k=6)보다 넓게 검색한다.
REPORT_TOP_K = 12

_SYSTEM = (
    "당신은 문서 분석 결과를 '구조화된 리포트'로 정리하는 도구입니다. "
    "반드시 아래 JSON 스키마로만 답하세요. 인사말·설명·마크다운·코드펜스 금지. "
    "오직 '참고 문맥'에 실제로 있는 내용만 사용하고, 문맥에 없는 수치는 절대 지어내지 마세요. "
    "수치 데이터가 문맥에 없으면 table 과 chart 를 null 로 두세요. "
    "모든 텍스트 값은 한국어로 작성하되, JSON 키는 스키마 그대로 영문을 쓰세요.\n"
    "스키마:\n"
    "{\n"
    '  "title": "리포트 제목",\n'
    '  "summary": "전체 요약 2~4문장",\n'
    '  "sections": [{"heading": "소제목", "paragraphs": ["문단", "문단"]}],\n'
    '  "table": {"columns": ["열1","열2"], "rows": [["값","값"]]} ,\n'
    '  "chart": {"type": "bar", "title": "차트제목", "categories": ["항목1","항목2"], '
    '"series": [{"name": "계열1", "values": [10, 20]}]}\n'
    "}\n"
    "규칙: chart.type 은 bar|line|pie 중 하나. chart 가 있으면 categories 길이와 각 "
    "series.values 길이가 같아야 합니다. 데이터가 없으면 table/chart 는 null."
)


def _chat_json(model: str, user_msg: str, num_ctx: int = rag.NUM_CTX) -> dict:
    """Ollama 채팅을 format=json 으로 호출해 dict 를 받는다. 파싱 실패 시 1회 재시도."""
    payload = {
        "model": model,
        "messages": [
            {"role": "system", "content": _SYSTEM},
            {"role": "user", "content": user_msg},
        ],
        "stream": False,
        "format": "json",
        "think": False,
        "options": {"num_ctx": num_ctx},
    }

    def _once(p: dict) -> str:
        r = requests.post(f"{OLLAMA_HOST}/api/chat", json=p, timeout=300)
        # 일부 모델은 think 파라미터 미지원 → 400. think 빼고 재시도.
        if r.status_code == 400 and "think" in p:
            p2 = dict(p)
            p2.pop("think", None)
            return _once(p2)
        r.raise_for_status()
        return (r.json().get("message") or {}).get("content", "")

    raw = _once(payload)
    data = _extract_json(raw)
    if data is None:
        # 안전망: 모델 로드 보장 후 1회 재시도.
        from backend import ollama as _ol
        _ol.load_model(model)
        raw = _once(payload)
        data = _extract_json(raw)
    if data is None:
        raise ValueError("모델이 올바른 JSON 리포트를 생성하지 못했습니다. 다시 시도해 주세요.")
    return data


def _extract_json(text: str) -> dict | None:
    """문자열에서 JSON 오브젝트를 추출. format=json 이면 보통 그대로 파싱된다."""
    if not text:
        return None
    text = text.strip()
    try:
        obj = json.loads(text)
        return obj if isinstance(obj, dict) else None
    except json.JSONDecodeError:
        pass
    # 코드펜스/잡텍스트가 섞인 경우 첫 { ~ 마지막 } 구간을 시도.
    s, e = text.find("{"), text.rfind("}")
    if s != -1 and e != -1 and e > s:
        try:
            obj = json.loads(text[s:e + 1])
            return obj if isinstance(obj, dict) else None
        except json.JSONDecodeError:
            return None
    return None


# ---------- content 모델 정규화/검증 ----------
def _as_text(v) -> str:
    return "" if v is None else str(v)


def _num(v):
    """셀 값을 숫자로 변환 가능하면 float, 아니면 None."""
    if isinstance(v, (int, float)):
        return float(v)
    if isinstance(v, str):
        m = re.search(r"-?\d[\d,]*\.?\d*", v)
        if m:
            try:
                return float(m.group(0).replace(",", ""))
            except ValueError:
                return None
    return None


def _clean_table(t) -> dict | None:
    """표 정규화: 사각형 보장, 빈 표는 None."""
    if not isinstance(t, dict):
        return None
    cols = [_as_text(c) for c in (t.get("columns") or [])]
    rows_in = t.get("rows") or []
    if not cols or not isinstance(rows_in, list):
        return None
    width = len(cols)
    rows = []
    for r in rows_in:
        if not isinstance(r, list):
            continue
        row = [_as_text(c) for c in r][:width]
        row += [""] * (width - len(row))   # 부족분 패딩
        rows.append(row)
    if not rows:
        return None
    return {"columns": cols, "rows": rows}


def _clean_chart(c) -> dict | None:
    """차트 정규화: type/계열/카테고리 길이 검증. 어긋나면 None(차트 생략)."""
    if not isinstance(c, dict):
        return None
    ctype = str(c.get("type") or "bar").lower()
    if ctype not in ("bar", "line", "pie"):
        ctype = "bar"
    cats = [_as_text(x) for x in (c.get("categories") or [])]
    series_in = c.get("series") or []
    if not cats or not isinstance(series_in, list):
        return None
    series = []
    for s in series_in:
        if not isinstance(s, dict):
            continue
        vals = s.get("values") or []
        nums = [_num(v) for v in vals]
        if any(n is None for n in nums) or len(nums) != len(cats):
            continue   # 숫자 검증·길이 불일치 계열은 버림(환각 방지)
        series.append({"name": _as_text(s.get("name")) or "값", "values": nums})
    if not series:
        return None
    # pie 는 첫 계열만 의미가 있다.
    if ctype == "pie":
        series = series[:1]
    return {"type": ctype, "title": _as_text(c.get("title")) or "차트",
            "categories": cats, "series": series}


def _clean_sections(s) -> list[dict]:
    out = []
    for sec in (s or []):
        if not isinstance(sec, dict):
            continue
        heading = _as_text(sec.get("heading")).strip()
        paras_in = sec.get("paragraphs") or []
        if isinstance(paras_in, str):
            paras_in = [paras_in]
        paras = [_as_text(p).strip() for p in paras_in if _as_text(p).strip()]
        if heading or paras:
            out.append({"heading": heading or "내용", "paragraphs": paras})
    return out


def build_report(question: str, doc_ids: list[str], model: str,
                 embed_model: str = rag.EMBED_MODEL) -> dict:
    """RAG 검색 + 구조화 LLM 호출로 포맷 무관 content 모델을 만든다."""
    if not model:
        raise ValueError("사용할 생성 모델이 선택되지 않았습니다.")
    chunks = rag.search(question, doc_ids, top_k=REPORT_TOP_K, embed_model=embed_model)
    if not chunks:
        raise ValueError("검색된 문서 내용이 없습니다. 먼저 문서를 추가하세요.")
    context, sources = rag.build_context(chunks)
    user_msg = f"참고 문맥:\n{context}\n\n요청: {question}\n\n위 문맥에 근거해 리포트 JSON 을 작성하세요."

    data = _chat_json(model, user_msg)

    content = {
        "title": _as_text(data.get("title")).strip() or (question[:40] or "분석 리포트"),
        "summary": _as_text(data.get("summary")).strip(),
        "sections": _clean_sections(data.get("sections")),
        "table": _clean_table(data.get("table")),
        "chart": _clean_chart(data.get("chart")),
        # 출처는 LLM 이 아니라 검색 결과(rag)에서 가져와 신뢰성 유지.
        "sources": _dedup_sources(sources),
        "meta": {
            "question": question,
            "model": model,
            "generated": time.strftime("%Y-%m-%d %H:%M"),
        },
    }
    return content


def _dedup_sources(sources: list[dict]) -> list[dict]:
    """같은 (문서, 페이지) 출처를 한 번만 남긴다."""
    seen, out = set(), []
    for s in sources:
        key = (s.get("name"), s.get("page"))
        if key in seen:
            continue
        seen.add(key)
        out.append({"name": s.get("name"), "page": s.get("page"), "doc_id": s.get("doc_id")})
    return out


def export(question: str, doc_ids: list[str], model: str, fmt: str, path: str,
           embed_model: str = rag.EMBED_MODEL) -> dict:
    """리포트를 만들고 fmt 에 맞는 파일로 저장. {ok, path, title} 반환."""
    fmt = (fmt or "").lower()
    if fmt not in FORMATS:
        raise ValueError(f"지원하지 않는 형식입니다: {fmt}")
    content = build_report(question, doc_ids, model, embed_model=embed_model)
    renderer = {
        "docx": to_docx, "pptx": to_pptx, "xlsx": to_xlsx, "txt": to_txt,
    }[fmt]
    renderer(content, path)
    return {"ok": True, "path": path, "title": content["title"]}


def _source_lines(content: dict) -> list[str]:
    out = []
    for s in content.get("sources") or []:
        loc = f" · p.{s['page']}" if s.get("page") else ""
        out.append(f"{s.get('name','')}{loc}")
    return out


# ==================== 렌더러: TXT ====================
def to_txt(content: dict, path: str) -> None:
    lines = [content["title"], "=" * len(content["title"]) , ""]
    meta = content.get("meta") or {}
    lines.append(f"생성: {meta.get('generated','')}  |  모델: {meta.get('model','')}")
    lines.append("")
    if content.get("summary"):
        lines += ["[ 요약 ]", content["summary"], ""]
    for sec in content.get("sections") or []:
        lines.append(f"■ {sec['heading']}")
        for p in sec["paragraphs"]:
            lines.append(p)
        lines.append("")
    table = content.get("table")
    if table:
        lines.append("[ 데이터 ]")
        lines += _txt_table(table["columns"], table["rows"])
        lines.append("")
    chart = content.get("chart")
    if chart:
        lines.append(f"[ 차트: {chart['title']} ({chart['type']}) ]")
        lines.append("카테고리: " + ", ".join(chart["categories"]))
        for s in chart["series"]:
            vals = ", ".join(_fmt_num(v) for v in s["values"])
            lines.append(f"  - {s['name']}: {vals}")
        lines.append("")
    srcs = _source_lines(content)
    if srcs:
        lines.append("[ 출처 ]")
        lines += [f"  - {s}" for s in srcs]
    with open(path, "w", encoding="utf-8") as f:
        f.write("\n".join(lines).rstrip() + "\n")


def _fmt_num(v) -> str:
    f = float(v)
    return str(int(f)) if f.is_integer() else f"{f:g}"


def _txt_table(columns: list[str], rows: list[list]) -> list[str]:
    """고정폭 텍스트 표."""
    widths = [len(str(c)) for c in columns]
    for r in rows:
        for i, c in enumerate(r):
            widths[i] = max(widths[i], len(str(c)))
    def fmt_row(cells):
        return " | ".join(str(c).ljust(widths[i]) for i, c in enumerate(cells))
    out = [fmt_row(columns), "-+-".join("-" * w for w in widths)]
    out += [fmt_row(r) for r in rows]
    return out


# ==================== 렌더러: DOCX ====================
def to_docx(content: dict, path: str) -> None:
    import docx
    from docx.shared import Pt, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    doc = docx.Document()
    title = doc.add_heading(content["title"], level=0)
    title.alignment = WD_ALIGN_PARAGRAPH.LEFT

    meta = content.get("meta") or {}
    sub = doc.add_paragraph()
    run = sub.add_run(f"생성: {meta.get('generated','')}    모델: {meta.get('model','')}")
    run.font.size = Pt(9)
    run.font.color.rgb = RGBColor(0x80, 0x80, 0x80)

    if content.get("summary"):
        doc.add_heading("요약", level=1)
        doc.add_paragraph(content["summary"])

    for sec in content.get("sections") or []:
        doc.add_heading(sec["heading"], level=1)
        for p in sec["paragraphs"]:
            doc.add_paragraph(p)

    table = content.get("table")
    if table:
        doc.add_heading("데이터", level=1)
        cols = table["columns"]
        t = doc.add_table(rows=1, cols=len(cols))
        t.style = "Table Grid"
        hdr = t.rows[0].cells
        for i, c in enumerate(cols):
            hdr[i].text = str(c)
            for para in hdr[i].paragraphs:
                for r in para.runs:
                    r.font.bold = True
        for row in table["rows"]:
            cells = t.add_row().cells
            for i, c in enumerate(row):
                cells[i].text = str(c)

    chart = content.get("chart")
    if chart:
        # docx 네이티브 차트는 까다로워, 데이터 표 + 캡션으로 싣는다.
        doc.add_heading(f"차트 데이터 · {chart['title']}", level=1)
        cats = chart["categories"]
        t = doc.add_table(rows=1, cols=len(cats) + 1)
        t.style = "Table Grid"
        hdr = t.rows[0].cells
        hdr[0].text = "계열"
        for i, cat in enumerate(cats):
            hdr[i + 1].text = str(cat)
        for s in chart["series"]:
            cells = t.add_row().cells
            cells[0].text = s["name"]
            for i, v in enumerate(s["values"]):
                cells[i + 1].text = _fmt_num(v)

    srcs = _source_lines(content)
    if srcs:
        doc.add_heading("출처", level=1)
        for s in srcs:
            doc.add_paragraph(s, style="List Bullet")

    doc.save(path)


# ==================== 렌더러: PPTX ====================
def to_pptx(content: dict, path: str) -> None:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE

    prs = Presentation()
    blank = prs.slide_layouts[6]      # 빈 레이아웃
    title_only = prs.slide_layouts[5]
    SW, SH = prs.slide_width, prs.slide_height

    def add_title_slide():
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = content["title"]
        meta = content.get("meta") or {}
        if slide.placeholders and len(slide.placeholders) > 1:
            slide.placeholders[1].text = f"{meta.get('generated','')}  ·  {meta.get('model','')}"

    def add_bullets_slide(heading, bullets):
        slide = prs.slides.add_slide(title_only)
        slide.shapes.title.text = heading
        box = slide.shapes.add_textbox(Inches(0.6), Inches(1.4),
                                       SW - Inches(1.2), SH - Inches(1.8))
        tf = box.text_frame
        tf.word_wrap = True
        for i, b in enumerate(bullets):
            para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            para.text = b
            para.font.size = Pt(16)
    # 표지
    add_title_slide()

    # 요약
    if content.get("summary"):
        add_bullets_slide("요약", [content["summary"]])

    # 섹션 (한 섹션 = 한 슬라이드, 문단을 불릿으로)
    for sec in content.get("sections") or []:
        add_bullets_slide(sec["heading"], sec["paragraphs"] or [""])

    # 차트 (네이티브)
    chart = content.get("chart")
    if chart:
        slide = prs.slides.add_slide(title_only)
        slide.shapes.title.text = chart["title"]
        cd = CategoryChartData()
        cd.categories = chart["categories"]
        for s in chart["series"]:
            cd.add_series(s["name"], s["values"])
        xl = {"bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
              "line": XL_CHART_TYPE.LINE_MARKERS,
              "pie": XL_CHART_TYPE.PIE}[chart["type"]]
        slide.shapes.add_chart(xl, Inches(0.8), Inches(1.4),
                               SW - Inches(1.6), SH - Inches(2.0), cd)

    # 데이터 표
    table = content.get("table")
    if table:
        _pptx_table_slide(prs, title_only, "데이터", table["columns"], table["rows"])

    # 출처
    srcs = _source_lines(content)
    if srcs:
        add_bullets_slide("출처", srcs)

    prs.save(path)


def _pptx_table_slide(prs, layout, heading, columns, rows):
    from pptx.util import Inches, Pt
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = heading
    n_rows = len(rows) + 1
    n_cols = len(columns)
    left, top = Inches(0.6), Inches(1.4)
    width = prs.slide_width - Inches(1.2)
    height = prs.slide_height - Inches(1.9)
    tbl = slide.shapes.add_table(n_rows, n_cols, left, top, width, height).table
    for i, c in enumerate(columns):
        cell = tbl.cell(0, i)
        cell.text = str(c)
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(12)
    for r, row in enumerate(rows, start=1):
        for i, c in enumerate(row):
            cell = tbl.cell(r, i)
            cell.text = str(c)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11)


# ==================== 렌더러: XLSX ====================
def to_xlsx(content: dict, path: str) -> None:
    from openpyxl import Workbook
    from openpyxl.chart import BarChart, LineChart, PieChart, Reference
    from openpyxl.styles import Font, Alignment

    wb = Workbook()
    bold = Font(bold=True)

    # --- 시트1: 요약 ---
    ws = wb.active
    ws.title = "요약"
    ws["A1"] = content["title"]
    ws["A1"].font = Font(bold=True, size=14)
    meta = content.get("meta") or {}
    ws["A2"] = f"생성: {meta.get('generated','')}   모델: {meta.get('model','')}"
    ws["A2"].font = Font(size=9, color="808080")
    r = 4
    if content.get("summary"):
        ws.cell(r, 1, "요약").font = bold
        r += 1
        cell = ws.cell(r, 1, content["summary"])
        cell.alignment = Alignment(wrap_text=True, vertical="top")
        r += 2
    for sec in content.get("sections") or []:
        ws.cell(r, 1, sec["heading"]).font = bold
        r += 1
        for p in sec["paragraphs"]:
            ws.cell(r, 1, p).alignment = Alignment(wrap_text=True, vertical="top")
            r += 1
        r += 1
    ws.column_dimensions["A"].width = 100

    # --- 시트2: 데이터 (+ 네이티브 차트) ---
    table = content.get("table")
    chart = content.get("chart")
    if table:
        wd = wb.create_sheet("데이터")
        for i, c in enumerate(table["columns"], start=1):
            wd.cell(1, i, c).font = bold
        for ri, row in enumerate(table["rows"], start=2):
            for ci, val in enumerate(row, start=1):
                num = _num(val)
                wd.cell(ri, ci, num if num is not None else val)
        _autofit(wd, table["columns"], table["rows"])

    if chart:
        wc = wb.create_sheet("차트")
        # 차트용 데이터 블록: 1행=카테고리 헤더 + 계열명, 이후 카테고리별 값.
        wc.cell(1, 1, "항목").font = bold
        for si, s in enumerate(chart["series"], start=2):
            wc.cell(1, si, s["name"]).font = bold
        for ri, cat in enumerate(chart["categories"], start=2):
            wc.cell(ri, 1, cat)
            for si, s in enumerate(chart["series"], start=2):
                wc.cell(ri, si, s["values"][ri - 2])
        n_cat = len(chart["categories"])
        n_ser = len(chart["series"])
        if chart["type"] == "pie":
            ch = PieChart()
            data = Reference(wc, min_col=2, min_row=1, max_row=1 + n_cat)
            cats = Reference(wc, min_col=1, min_row=2, max_row=1 + n_cat)
            ch.add_data(data, titles_from_data=True)
            ch.set_categories(cats)
        else:
            ch = BarChart() if chart["type"] == "bar" else LineChart()
            data = Reference(wc, min_col=2, max_col=1 + n_ser, min_row=1, max_row=1 + n_cat)
            cats = Reference(wc, min_col=1, min_row=2, max_row=1 + n_cat)
            ch.add_data(data, titles_from_data=True)
            ch.set_categories(cats)
        ch.title = chart["title"]
        ch.height = 9
        ch.width = 16
        wc.add_chart(ch, "E2")

    # --- 시트3: 출처 ---
    srcs = _source_lines(content)
    if srcs:
        wsrc = wb.create_sheet("출처")
        wsrc.cell(1, 1, "출처").font = bold
        for i, s in enumerate(srcs, start=2):
            wsrc.cell(i, 1, s)
        wsrc.column_dimensions["A"].width = 60

    wb.save(path)


def _autofit(ws, columns, rows) -> None:
    from openpyxl.utils import get_column_letter
    for i, col in enumerate(columns, start=1):
        width = len(str(col))
        for row in rows:
            if i - 1 < len(row):
                width = max(width, len(str(row[i - 1])))
        ws.column_dimensions[get_column_letter(i)].width = min(max(width + 2, 10), 50)
